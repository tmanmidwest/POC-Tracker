"""Server-side PowerPoint (.pptx) rendering of a single-POC executive readout.

Built with python-pptx (pure Python — no system libraries, same rationale as the
Word path in ``report_docx``). This is a *presentation* view of a POC, not the
full report: a title slide, an optional executive summary, a use-case scorecard
(with a doughnut progress ring and pass/fail chips), per-category result slides,
screenshot "proof" slides, and a next-steps slide.

Deliberately deterministic: every number comes from the project's own data (the
same ``context`` dict the PDF/Word/on-screen report use), so the deck can be
generated with no AI provider configured. When a ``ReadoutNarrative`` is passed
in, its AI-written bullets enrich the summary/next-steps slides; the layout never
depends on it. Speaker notes are written for every slide so an AE can present
straight from the deck.

Branding:
* ``template_path`` — an admin-supplied .pptx whose slide master/theme/fonts sit
  underneath the generated slides. The deck is laid out on a base 16:9 design but
  every coordinate is scaled to the template's *actual* slide size (see Canvas),
  so a 4:3 (or any) template is honored rather than forced to 16:9.
* ``logo_path`` — an image stamped in the top-right corner of every slide.
"""

from __future__ import annotations

import io
import logging
from dataclasses import dataclass, field
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from app.models import Project
from app.services import screenshots as screenshot_store

log = logging.getLogger(__name__)

# The deck is authored against this base 16:9 design; Canvas scales it to fit the
# real slide size (which may come from a 4:3 template, etc.).
_BASE_W_IN = 13.333
_BASE_H_IN = 7.5
_MARGIN_IN = 0.6

# Light palette (default): dark text on the template's light layout.
_MUTED = RGBColor(0x64, 0x74, 0x8B)
_INK = RGBColor(0x1E, 0x29, 0x3B)
_GREEN = RGBColor(0x16, 0xA3, 0x4A)
_AMBER = RGBColor(0xB4, 0x54, 0x09)
_TRACK = RGBColor(0xE2, 0xE8, 0xF0)  # light ring "remaining" track
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Dark palette: light text/muted + a dark "remaining" track. _DARK_BG is a fallback
# slide background painted only when the chosen layout isn't itself dark. Status
# colors (green/amber) and the brand accent read fine on either palette, so they
# stay shared.
_DARK_INK = RGBColor(0xF1, 0xF5, 0xF9)
_DARK_MUTED = RGBColor(0x94, 0xA3, 0xB8)
_DARK_TRACK = RGBColor(0x33, 0x41, 0x55)
_DARK_BG = RGBColor(0x0F, 0x17, 0x2A)

# Keep decks reasonable: at most this many use cases listed per category slide
# (more spill onto a continuation slide) and this many screenshot slides total.
_UC_PER_SLIDE = 9
_MAX_SCREENSHOTS = 12


class Canvas:
    """Maps the base 16:9 design onto a slide of arbitrary size.

    ``x``/``w`` scale by the width ratio, ``y``/``h`` by the height ratio, and
    fonts by the smaller of the two so text never overflows a narrower canvas.
    On a 16:9 slide the ratios are 1.0 and everything is pixel-identical to the
    original design; on a 4:3 template widths shrink to ~0.75 and text with them.
    """

    def __init__(self, emu_w: int, emu_h: int) -> None:
        self.W = int(emu_w)
        self.H = int(emu_h)
        self._sx = self.W / Inches(_BASE_W_IN)
        self._sy = self.H / Inches(_BASE_H_IN)
        self._sf = min(self._sx, self._sy)
        self.MARGIN = self.x(_MARGIN_IN)
        self.CW = self.W - 2 * self.MARGIN  # content width inside the margins

    def x(self, in_: float) -> int:
        return int(Inches(in_) * self._sx)

    def y(self, in_: float) -> int:
        return int(Inches(in_) * self._sy)

    def w(self, in_: float) -> int:
        return int(Inches(in_) * self._sx)

    def h(self, in_: float) -> int:
        return int(Inches(in_) * self._sy)

    def pt(self, size: float) -> Pt:
        return Pt(round(size * self._sf, 1))


@dataclass
class _Deck:
    """Per-build context threaded through the slide builders."""

    prs: Presentation
    cv: Canvas
    accent: RGBColor
    # Theme palette (defaults to light). ``dark`` also steers layout selection and
    # the fallback background in _new_slide.
    ink: RGBColor = _INK
    muted: RGBColor = _MUTED
    track: RGBColor = _TRACK
    bg: RGBColor = _DARK_BG
    dark: bool = False
    logo_path: str | None = None
    # The customer's own logo, stamped on the title slide (co-branding).
    customer_logo_path: str | None = None
    # When a template is in play, the template carries the branding — we don't
    # stamp the app's brand name onto the deck.
    used_template: bool = False


@dataclass
class ReadoutNarrative:
    """Optional AI-written, slide-ready copy for the deck (see ai.readout)."""

    summary: list[str] = field(default_factory=list)      # exec-summary bullets
    next_steps: list[str] = field(default_factory=list)   # recommended next steps
    model_label: str = ""                                  # "anthropic/claude-…"

    def is_empty(self) -> bool:
        return not (self.summary or self.next_steps)


def _accent(color: str | None) -> RGBColor:
    """Brand accent as an RGBColor, mirroring report_docx (empty -> theme ink)."""
    raw = (color or "#1e293b").lstrip("#")
    if len(raw) != 6:
        raw = "1e293b"
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def _template_accent(prs: Presentation) -> str | None:
    """The template theme's primary accent (accent1) as a hex string, if any.

    This is what a branded deck follows — the app's own Branding color never
    touches a templated output. Note the value is whatever the template declares:
    if a template puts an odd color in accent1, the deck honors it (fix the
    template to change it).
    """
    try:
        from lxml import etree

        master = prs.slide_masters[0]
        theme_part = next(
            (rel.target_part for rel in master.part.rels.values() if "theme" in rel.reltype),
            None,
        )
        if theme_part is None:
            return None
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        root = etree.fromstring(theme_part.blob)
        srgb = root.find(".//a:clrScheme/a:accent1/a:srgbClr", ns)
        if srgb is not None:
            return srgb.get("val")
        sysclr = root.find(".//a:clrScheme/a:accent1/a:sysClr", ns)
        if sysclr is not None:
            return sysclr.get("lastClr")
    except Exception:
        return None
    return None


def _blank_layout(prs: Presentation, dark: bool = False):
    """Pick a blank layout to build on, matching the requested theme.

    Templates commonly ship both a dark and a light "Blank" layout. Light content
    is dark-text-on-light and dark content is light-text-on-dark, so each theme
    prefers its own blank (avoiding gradients), falling back through progressively
    looser matches to python-pptx's default "Blank" at index 6. In dark mode with no
    dark layout available, _new_slide paints a dark background instead.
    """
    layouts = list(prs.slide_layouts)

    def first(pred):
        return next((lay for lay in layouts if lay.name and pred(lay.name.lower())), None)

    if dark:
        return (
            first(lambda n: "blank" in n and "dark" in n and "gradient" not in n)
            or first(lambda n: "blank" in n and "dark" in n)
            or first(lambda n: "blank" in n)
            or (layouts[6] if len(layouts) > 6 else layouts[-1])
        )
    return (
        first(lambda n: "blank" in n and "light" in n)
        or first(lambda n: "blank" in n and "dark" not in n and "gradient" not in n)
        or first(lambda n: "blank" in n and "dark" not in n)
        or first(lambda n: "blank" in n)
        or (layouts[6] if len(layouts) > 6 else layouts[-1])
    )


def _clear_existing_slides(prs: Presentation) -> None:
    """Drop any slides a template ships with, keeping its masters/layouts/theme.

    Corporate templates commonly include example/boilerplate slides; without this
    they'd land in front of the generated readout. Removing the ``sldId`` entries
    (and their relationships) leaves the theme and layouts intact.
    """
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        rid = sld_id.get(qn("r:id"))
        if rid:
            try:
                prs.part.drop_rel(rid)
            except KeyError:
                pass
        sld_id_lst.remove(sld_id)


def _new_slide(d: _Deck):
    """Add a blank slide and stamp the brand logo (if any) on it."""
    layout = _blank_layout(d.prs, dark=d.dark)
    slide = d.prs.slides.add_slide(layout)
    # In dark mode, if the chosen layout isn't itself dark (no template, or a
    # template without a dark blank), lay down a full-bleed dark background first
    # so it sits behind all content.
    if d.dark and not (layout.name and "dark" in layout.name.lower()):
        _rect(slide, 0, 0, d.cv.W, d.cv.H, d.bg)
    if d.logo_path:
        _stamp_logo(slide, d.cv, d.logo_path)
    return slide


def _stamp_logo(slide, cv: Canvas, logo_path: str) -> None:
    """Place the brand logo in the top-right corner, scaled to a small height."""
    try:
        pic = slide.shapes.add_picture(logo_path, 0, 0, height=cv.h(0.5))
    except Exception:  # unreadable/unsupported image — just skip it
        log.warning("pptx_logo_skipped", extra={"logo_path": logo_path})
        return
    pic.left = int(cv.W - cv.x(0.35) - pic.width)
    pic.top = cv.y(0.28)


def _notes(slide, text: str) -> None:
    """Fill the slide's speaker-notes pane (talking points for the presenter)."""
    if text:
        slide.notes_slide.notes_text_frame.text = text.strip()


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _run(paragraph, text, *, size, color=_INK, bold=False):
    r = paragraph.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    return r


def _rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _accent_bar(d: _Deck, slide, *, top, height=None):
    """A thin full-bleed accent rule used to brand each slide."""
    cv = d.cv
    return _rect(slide, 0, top, cv.W, height if height is not None else cv.h(0.12), d.accent)


def _slide_heading(d: _Deck, slide, kicker, title):
    """Standard content-slide header: small colored kicker + big title + rule."""
    cv = d.cv
    _, tf = _textbox(slide, cv.MARGIN, cv.y(0.45), cv.CW, cv.h(1.1))
    p1 = tf.paragraphs[0]
    _run(p1, kicker.upper(), size=cv.pt(12), bold=True, color=d.accent)
    p2 = tf.add_paragraph()
    _run(p2, title, size=cv.pt(30), bold=True, color=d.ink)
    _accent_bar(d, slide, top=cv.y(1.55), height=cv.h(0.04))


def _uc_complete(uc) -> bool:
    return bool(uc.status and uc.status.is_complete_status)


def _uc_status_bits(uc):
    """(glyph, label, color) for a pass/fail-style status chip."""
    if _uc_complete(uc):
        return "✓", "PASS", _GREEN
    label = uc.status.name if uc.status else "Not started"
    return "◔", label, _AMBER


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------


def _title_slide(d: _Deck, project, ctx):
    cv = d.cv
    slide = _new_slide(d)

    # Customer logo, top-left above the title (co-branding). Best-effort.
    if d.customer_logo_path:
        try:
            slide.shapes.add_picture(
                d.customer_logo_path, cv.MARGIN, cv.y(0.95), height=cv.h(0.85)
            )
        except Exception:
            log.warning(
                "pptx_customer_logo_skipped",
                extra={"path": d.customer_logo_path},
            )

    _accent_bar(d, slide, top=cv.y(2.55))

    _, tf = _textbox(slide, cv.MARGIN, cv.y(2.75), cv.CW, cv.h(2.6))

    # The app brand name is app chrome — only stamp it on the built-in deck. With a
    # template, the template's own master/logo carries the branding.
    have_brand_line = False
    if not d.used_template:
        brand = ctx["branding"].get("name") or "Questlog"
        _run(tf.paragraphs[0], brand, size=cv.pt(14), bold=True, color=d.accent)
        have_brand_line = True

    p_title = tf.add_paragraph() if have_brand_line else tf.paragraphs[0]
    if have_brand_line:
        p_title.space_before = cv.pt(6)
    _run(p_title, project.display_name, size=cv.pt(40), bold=True, color=d.ink)

    status = project.status.name if project.status else "—"
    customer = project.customer.name if project.customer else "—"
    sub = f"{customer}  ·  {status}"
    if project.is_archived:
        sub += "  ·  Archived"
    p_sub = tf.add_paragraph()
    p_sub.space_before = cv.pt(10)
    _run(p_sub, sub, size=cv.pt(18), color=d.muted)

    p_meta = tf.add_paragraph()
    p_meta.space_before = cv.pt(4)
    _run(p_meta, f"Executive readout · Generated {ctx['generated_on']}",
         size=cv.pt(12), color=d.muted)

    se = project.sales_engineer.display_label if project.sales_engineer else "—"
    _notes(
        slide,
        f"Executive readout for {project.display_name} ({customer}). "
        f"Current status: {status}. Sales Engineer: {se}. "
        f"Account Executive: {project.account_executive or '—'}. "
        "Open by framing the goal of this POC and who's in the room.",
    )


def _summary_slide(d: _Deck, project, narrative):
    cv = d.cv
    bullets = narrative.summary if narrative else []
    if not bullets and not project.exec_summary:
        return
    slide = _new_slide(d)
    _slide_heading(d, slide, "Executive summary", project.display_name)
    _, tf = _textbox(slide, cv.MARGIN, cv.y(1.85), cv.CW, cv.h(5.0))
    first = True
    if bullets:
        for b in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_after = cv.pt(9)
            _run(p, "•  ", size=cv.pt(16), bold=True, color=d.accent)
            _run(p, b.strip(), size=cv.pt(16), color=d.ink)
            first = False
        _notes(slide, "Lead with the outcome, then the two or three points that matter "
                      "most to this audience. Pause for reactions before diving into detail.")
    else:
        for para in project.exec_summary.split("\n\n"):
            para = para.strip()
            if not para:
                continue
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_after = cv.pt(10)
            _run(p, para, size=cv.pt(16), color=d.ink)
            first = False
        _notes(slide, "Walk through the executive summary at a high level; don't read it "
                      "verbatim. Land the headline outcome first.")


def _progress_ring(d: _Deck, slide, done, total, *, left, top, size):
    """A doughnut chart used as a progress ring, with the % in the middle."""
    cv = d.cv
    remaining = max(total - done, 0)
    data = CategoryChartData()
    data.categories = ["Validated", "Remaining"]
    data.add_series("Progress", (done, remaining))
    frame = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, left, top, size, size, data)
    chart = frame.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.has_data_labels = False
    pts = list(plot.series[0].points)
    pts[0].format.fill.solid()
    pts[0].format.fill.fore_color.rgb = d.accent
    pts[1].format.fill.solid()
    pts[1].format.fill.fore_color.rgb = d.track
    for pt in pts:  # hide the thin segment borders
        pt.format.line.fill.background()
    dough = plot._element
    hole = dough.find(qn("c:holeSize"))
    if hole is None:
        hole = dough.makeelement(qn("c:holeSize"), {})
        dough.append(hole)
    hole.set("val", "62")

    pct = round(done / total * 100) if total else 0
    _, tf = _textbox(slide, left, top + Emu(int(size * 0.34)), size, Emu(int(size * 0.32)))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, f"{pct}%", size=cv.pt(36), bold=True, color=d.accent)


def _scorecard_slide(d: _Deck, project, ctx):
    cv = d.cv
    slide = _new_slide(d)
    _slide_heading(d, slide, "Results at a glance", "Use-case scorecard")
    progress = ctx["progress"]
    total = progress["total"]

    if total:
        _progress_ring(
            d, slide, progress["done"], total,
            left=cv.MARGIN, top=cv.y(2.1), size=cv.h(3.4),
        )
    _, tf = _textbox(slide, cv.MARGIN, cv.y(5.6), cv.w(3.4), cv.h(1.2))
    cap = tf.paragraphs[0]
    cap.alignment = PP_ALIGN.CENTER
    _run(cap, f"{progress['done']} of {total} ", size=cv.pt(20), bold=True, color=d.ink)
    _run(cap, "validated", size=cv.pt(20), color=d.muted)

    # Status breakdown as pass/fail-colored rows on the right.
    counts: dict[str, tuple[int, RGBColor]] = {}
    for uc in project.use_cases:
        label = uc.status.name if uc.status else "—"
        color = _GREEN if _uc_complete(uc) else _AMBER
        n, _ = counts.get(label, (0, color))
        counts[label] = (n + 1, color)
    _, tf2 = _textbox(slide, cv.x(5.0), cv.y(2.1), cv.w(7.7), cv.h(4.6))
    head = tf2.paragraphs[0]
    _run(head, "BY STATUS", size=cv.pt(12), bold=True, color=d.muted)
    for label, (n, color) in sorted(counts.items(), key=lambda kv: (-kv[1][0], kv[0].lower())):
        row = tf2.add_paragraph()
        row.space_after = cv.pt(6)
        _run(row, f"{n}   ", size=cv.pt(20), bold=True, color=color)
        _run(row, label, size=cv.pt(18), color=d.ink)
    if not counts:
        _run(tf2.add_paragraph(), "No use cases yet.", size=cv.pt(16), color=d.muted)

    _notes(
        slide,
        f"Headline number: {progress['done']} of {total} use cases validated "
        f"({progress['pct']}%). Use the ring to anchor the conversation, then use the "
        "status breakdown to set up which categories you'll walk through next.",
    )


def _category_slides(d: _Deck, ctx):
    cv = d.cv
    for group in ctx["use_case_groups"]:
        ucs = group["use_cases"]
        chunks = [ucs[i : i + _UC_PER_SLIDE] for i in range(0, len(ucs), _UC_PER_SLIDE)] or [[]]
        for idx, chunk in enumerate(chunks):
            slide = _new_slide(d)
            title = group["category"]
            if len(chunks) > 1:
                title = f"{title} ({idx + 1}/{len(chunks)})"
            _slide_heading(d, slide, "Use cases", title)
            _, tf = _textbox(slide, cv.MARGIN, cv.y(1.85), cv.CW, cv.h(5.1))
            first = True
            for uc in chunk:
                glyph, label, color = _uc_status_bits(uc)
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.space_after = cv.pt(7)
                first = False
                name = uc.name
                if uc.reference_number:
                    name = f"{uc.reference_number}  {name}"
                _run(p, f"{glyph} ", size=cv.pt(15), bold=True, color=color)
                _run(p, name, size=cv.pt(15), bold=True, color=d.ink)
                _run(p, f"   {label}", size=cv.pt(13), bold=True, color=color)
                if uc.comments:
                    note = uc.comments.strip().replace("\n", " ")
                    if len(note) > 160:
                        note = note[:157] + "…"
                    sub = tf.add_paragraph()
                    sub.space_after = cv.pt(7)
                    sub.level = 1
                    _run(sub, note, size=cv.pt(12), color=d.muted)
            passed = sum(1 for uc in chunk if _uc_complete(uc))
            _notes(
                slide,
                f"{group['category']}: {passed} of {len(chunk)} shown are validated. "
                "Call out the wins first, then be candid about anything still open and "
                "what it will take to close it.",
            )


def _screenshot_slides(d: _Deck, ctx):
    shown = 0
    for group in ctx["use_case_groups"]:
        for uc in group["use_cases"]:
            for shot in uc.screenshots or []:
                if shown >= _MAX_SCREENSHOTS:
                    return
                try:
                    path = str(screenshot_store.path_for(shot))
                except Exception:
                    continue
                slide = _new_slide(d)
                name = uc.name
                if uc.reference_number:
                    name = f"{uc.reference_number}  {name}"
                _slide_heading(d, slide, "Proof", name)
                _place_image(d, slide, path, shot)
                _notes(slide, f"Evidence for “{uc.name}”. "
                              f"{shot.caption or 'Point out what the customer is seeing and why it matters.'}")
                shown += 1


def _place_image(d: _Deck, slide, path, shot) -> None:
    """Fit an image inside the content area, centered, preserving aspect ratio."""
    cv = d.cv
    area_top = cv.y(1.9)
    caption_h = cv.h(0.55) if shot.caption else cv.h(0.1)
    max_w = cv.CW
    max_h = cv.H - area_top - caption_h - cv.h(0.3)
    try:
        pic = slide.shapes.add_picture(path, cv.MARGIN, area_top, width=max_w)
    except Exception:
        log.warning("pptx_screenshot_skipped", extra={"screenshot_id": shot.id})
        _, tf = _textbox(slide, cv.MARGIN, area_top, cv.CW, cv.h(1.0))
        _run(tf.paragraphs[0], "[screenshot unavailable]", size=cv.pt(14), color=d.muted)
        return
    if pic.height > max_h:  # too tall for the slot — constrain by height instead
        pic.width = int(pic.width * (max_h / pic.height))
        pic.height = int(max_h)
    pic.left = int((cv.W - pic.width) / 2)
    if shot.caption:
        _, tf = _textbox(
            slide, cv.MARGIN, Emu(int(pic.top) + int(pic.height)) + cv.h(0.1),
            cv.CW, caption_h,
        )
        cap = tf.paragraphs[0]
        cap.alignment = PP_ALIGN.CENTER
        _run(cap, shot.caption, size=cv.pt(12), color=d.muted)


def _next_steps_slide(d: _Deck, project, ctx, narrative):
    cv = d.cv
    slide = _new_slide(d)
    _slide_heading(d, slide, "Where we are", "Timeline & next steps")

    def _fmt(dt):
        return dt.strftime("%b %-d, %Y") if dt else "—"

    open_ucs = [
        uc
        for g in ctx["use_case_groups"]
        for uc in g["use_cases"]
        if not _uc_complete(uc)
    ]

    # Left column: timeline + open items (deterministic).
    _, tf = _textbox(slide, cv.MARGIN, cv.y(1.85), cv.w(6.0), cv.h(5.1))
    p = tf.paragraphs[0]
    _run(p, "Timeline", size=cv.pt(16), bold=True, color=d.accent)
    p2 = tf.add_paragraph()
    _run(p2, f"{_fmt(project.start_date)}  →  {_fmt(project.end_date)}",
         size=cv.pt(16), color=d.ink)

    head = tf.add_paragraph()
    head.space_before = cv.pt(14)
    _run(head, f"Open items ({len(open_ucs)})", size=cv.pt(16), bold=True, color=d.accent)
    if not open_ucs:
        _run(tf.add_paragraph(), "All tracked use cases are complete.",
             size=cv.pt(15), color=_GREEN)
    for uc in open_ucs[:10]:
        row = tf.add_paragraph()
        row.space_after = cv.pt(4)
        name = uc.name
        if uc.reference_number:
            name = f"{uc.reference_number}  {name}"
        _run(row, "○  ", size=cv.pt(14), bold=True, color=_AMBER)
        _run(row, name, size=cv.pt(14), color=d.ink)
    if len(open_ucs) > 10:
        _run(tf.add_paragraph(), f"…and {len(open_ucs) - 10} more",
             size=cv.pt(12), color=d.muted)

    # Right column: AI-recommended next steps, when available.
    steps = narrative.next_steps if narrative else []
    if steps:
        _, tf2 = _textbox(slide, cv.x(6.9), cv.y(1.85), cv.w(5.8), cv.h(5.1))
        h = tf2.paragraphs[0]
        _run(h, "Recommended next steps", size=cv.pt(16), bold=True, color=d.accent)
        for s in steps:
            row = tf2.add_paragraph()
            row.space_before = cv.pt(6)
            _run(row, "→  ", size=cv.pt(15), bold=True, color=d.accent)
            _run(row, s.strip(), size=cv.pt(15), color=d.ink)

    note = (
        f"{len(open_ucs)} open item(s). Close by confirming the timeline, naming the "
        "owner for each open item, and agreeing on the decision date."
    )
    if steps:
        note += " Use the recommended next steps to drive to a mutual action plan."
    _notes(slide, note)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------


def build_project_readout_pptx(
    project: Project,
    context: dict[str, Any],
    *,
    narrative: ReadoutNarrative | None = None,
    template_path: str | None = None,
    logo_path: str | None = None,
    customer_logo_path: str | None = None,
    theme: str = "light",
) -> bytes:
    """Build an executive readout deck for ``project`` from its report context.

    ``context`` is the same dict produced for the PDF/Word report (branding,
    progress, use_case_groups, generated_on). ``narrative`` (optional) supplies
    AI-written bullets. ``template_path`` (optional) is a branded .pptx whose
    master/theme underlies the slides — its own slide size is honored. ``logo_path``
    (optional) is stamped in the top-right of every slide. ``theme`` is ``"light"``
    (default) or ``"dark"``: dark uses the template's dark blank layout (or a painted
    dark background) with a light text palette. Returns the .pptx bytes.
    """
    dark = str(theme).lower() == "dark"
    if narrative is not None and narrative.is_empty():
        narrative = None

    used_template = False
    if template_path:
        try:
            prs = Presentation(template_path)
            used_template = True
        except Exception:  # unreadable/corrupt template — fall back to the built-in
            log.warning("pptx_template_unreadable", extra={"template_path": template_path})
            prs = Presentation()
    else:
        prs = Presentation()

    if used_template:
        # A template may ship with example slides — drop them, keep its theme, and
        # lay out relative to its authored slide size.
        _clear_existing_slides(prs)
    else:
        # No template: use our own 16:9 canvas.
        prs.slide_width = Inches(_BASE_W_IN)
        prs.slide_height = Inches(_BASE_H_IN)

    # Accent source is a hard split by whether a template is in play:
    #  * templated output → the TEMPLATE's own theme accent; the app's Branding
    #    color is intentionally ignored so app chrome never leaks into a branded
    #    deck (falls back to the neutral slate default if the theme has no accent).
    #  * built-in deck (no template) → the app's Branding accent, since that's the
    #    only branding signal available.
    if used_template:
        accent = _accent(_template_accent(prs))
    else:
        accent = _accent(context["branding"].get("color"))

    cv = Canvas(prs.slide_width, prs.slide_height)
    deck = _Deck(
        prs=prs, cv=cv, accent=accent,
        ink=_DARK_INK if dark else _INK,
        muted=_DARK_MUTED if dark else _MUTED,
        track=_DARK_TRACK if dark else _TRACK,
        dark=dark,
        logo_path=logo_path, customer_logo_path=customer_logo_path,
        used_template=used_template,
    )

    _title_slide(deck, project, context)
    _summary_slide(deck, project, narrative)
    _scorecard_slide(deck, project, context)
    _category_slides(deck, context)
    _screenshot_slides(deck, context)
    _next_steps_slide(deck, project, context, narrative)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
