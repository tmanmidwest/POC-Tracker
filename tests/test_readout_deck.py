"""Executive readout deck: AI narrative path and the JSON parser (provider mocked)."""

from __future__ import annotations

import io
import zipfile

import pytest
from fastapi.testclient import TestClient

from app.db import get_session_factory
from app.models import AIProvider, Customer, Project, ProjectStatus


@pytest.fixture
def admin_ui(client: TestClient) -> TestClient:
    from app.config import get_settings

    s = get_settings()
    resp = client.post(
        "/ui/login",
        data={"username": s.initial_admin_username, "password": s.initial_admin_password},
        follow_redirects=False,
    )
    assert resp.status_code == 303, resp.text
    return client


def _make_project(name: str) -> int:
    db = get_session_factory()()
    try:
        customer = Customer(name=f"Cust {name}")
        db.add(customer)
        db.flush()
        status = db.query(ProjectStatus).order_by(ProjectStatus.sort_order).first()
        project = Project(customer_id=customer.id, name=name, status_id=status.id)
        db.add(project)
        db.commit()
        return project.id
    finally:
        db.close()


def _add_provider() -> None:
    from app.services.secret_box import encrypt_secret

    db = get_session_factory()()
    try:
        db.add(
            AIProvider(
                provider="anthropic", display_name="Claude", model="claude-opus-4-8",
                api_key_encrypted=encrypt_secret("sk-test"),
                is_enabled=True, is_default=True,
            )
        )
        db.commit()
    finally:
        db.close()


def _patch_generate(monkeypatch, fn) -> None:
    from app.services.ai import registry

    monkeypatch.setitem(
        registry.PROVIDERS,
        "anthropic",
        registry.PROVIDERS["anthropic"].__class__(
            key="anthropic", label="Anthropic (Claude)",
            default_model="claude-opus-4-8", suggested_models=["claude-opus-4-8"],
            implemented=True, generate=fn,
        ),
    )


def _slide_text(pptx_bytes: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(pptx_bytes))
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return "\n".join(out)


def _run_colors(pptx_bytes: bytes) -> set[str]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(pptx_bytes))
    colors: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color and run.font.color.rgb is not None:
                            colors.add(str(run.font.color.rgb))
                    except (AttributeError, TypeError):
                        pass
    return colors


def test_readout_theme_light_vs_dark(admin_ui: TestClient) -> None:
    """theme=dark flips the text palette to light ink; light keeps dark ink."""
    pid = _make_project("Theme POC")
    light = admin_ui.get(f"/ui/reports/projects/{pid}/readout.pptx?theme=light")
    dark = admin_ui.get(f"/ui/reports/projects/{pid}/readout.pptx?theme=dark")
    assert light.status_code == 200 and dark.status_code == 200

    light_colors = _run_colors(light.content)
    dark_colors = _run_colors(dark.content)
    assert "1E293B" in light_colors        # light deck uses dark ink for body text
    assert "F1F5F9" in dark_colors         # dark deck uses light ink for body text
    assert "F1F5F9" not in light_colors    # light ink is dark-mode-only


def test_ai_deck_renders_generated_bullets(
    admin_ui: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    pid = _make_project("Narrative POC")
    _add_provider()

    def fake_generate(*, api_key, model, system, prompt, max_tokens=1500, usage=None):
        return (
            'Here you go: {"summary": ["Core scenarios validated in staging", '
            '"Strong champion in security team"], '
            '"next_steps": ["Schedule executive readout", "Confirm procurement path"]}'
        )

    _patch_generate(monkeypatch, fake_generate)

    r = admin_ui.get(f"/ui/reports/projects/{pid}/readout.pptx?ai=1")
    assert r.status_code == 200
    text = _slide_text(r.content)
    assert "Core scenarios validated in staging" in text
    assert "Schedule executive readout" in text


def test_ai_deck_falls_back_on_bad_json(
    admin_ui: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    pid = _make_project("Fallback POC")
    _add_provider()

    def fake_generate(*, api_key, model, system, prompt, max_tokens=1500, usage=None):
        return "sorry, I can't help with that"  # no JSON

    _patch_generate(monkeypatch, fake_generate)

    # Deck must still be produced (deterministic content), not a 500.
    r = admin_ui.get(f"/ui/reports/projects/{pid}/readout.pptx?ai=1")
    assert r.status_code == 200
    names = zipfile.ZipFile(io.BytesIO(r.content)).namelist()
    assert any(n.startswith("ppt/slides/slide") for n in names), names


def test_narrative_parser_handles_junk() -> None:
    from app.services.ai.readout import _parse

    assert _parse("no json here") is None
    assert _parse("{not valid json}") is None
    # Extracts a JSON object embedded in prose and caps/cleans the lists.
    summary, steps = _parse('prefix {"summary": ["a", "", 3, "b"], "next_steps": ["x"]} suffix')
    assert summary == ["a", "b"]
    assert steps == ["x"]
